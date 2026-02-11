import asyncio
import logging
from pathlib import Path
from typing import List, Dict, Any, TypedDict, Tuple, Optional, NotRequired
from langgraph.graph import StateGraph, END
from langchain_openai import ChatOpenAI
from langchain_core.prompts import ChatPromptTemplate
from langchain_core.runnables import RunnableConfig
from pptx import Presentation
from pptx.dml.color import RGBColor
from collections import defaultdict
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
import re
import shutil
import tempfile
import time
from statistics import median

from utils import *

# 定义全局的 logger
logger = logging.getLogger(__name__)


# ==========================================
# 1. 定义 Agent State
# ==========================================
class AgentState(TypedDict):
    input_ppt_path: str        
    output_ppt_path: NotRequired[str]
    target_language: str = "English"
    extracted_data: NotRequired[List[Dict]]
    translation_map: NotRequired[Dict]    
    status_msg: NotRequired[str]
    max_concurrent: NotRequired[int]
    batch_size: NotRequired[int]

# ==========================================
# 1. 节点一：解析PPT并提取文本
# ==========================================

def node_parse_ppt(state: AgentState) -> AgentState:
    """同步节点：解析 PPT 并提取文本"""
    logger.info("🔍 开始解析 PPT...")
    prs = Presentation(state['input_ppt_path'])
    extracted_data = []
    
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            text = shape.text.strip()
            if not text:
                continue
            
            extracted_data.append({
                "slide_index": slide_idx,
                "original_text": text, 
            })
    
    state["extracted_data"] = extracted_data
    state["status_msg"] = f"✅ 解析完成：提取了 {len(extracted_data)} 个文本块"
    logger.info(f"📊 解析完成：{len(extracted_data)} 个文本块")
    return state

# ==========================================
# 2. 节点二：使用异步 LLM 进行高效的并发翻译
# ==========================================

async def async_node_translate_text(llm, state: AgentState) -> AgentState:
    """
    异步节点：使用异步 LLM 进行高效的并发翻译
    """
    logger.info("🌍 开始翻译...")
    translation_instruction = load_prompt("./prompts/translation_instruction.txt")
    prompt = ChatPromptTemplate.from_messages([
        ("system", translation_instruction),
        ("user", "{text}")
    ])
    chain = prompt | llm
    
    translation_map = {}
    batch_texts = [(item["original_text"], item) for item in state["extracted_data"]]
    
    # 并发控制参数
    MAX_CONCURRENT = state.get('max_concurrent', 10)
    BATCH_SIZE = state.get('batch_size', 10)
    logger.info(f"配置: 并发数={MAX_CONCURRENT}, 批次大小={BATCH_SIZE}")

    MAX_RETRIES = 2
    
    semaphore = asyncio.Semaphore(MAX_CONCURRENT)
    
    async def translate_single(text: str) -> Tuple[str, Optional[str]]:
        """翻译单个文本，带重试"""
        async with semaphore:
            for attempt in range(MAX_RETRIES + 1):
                try:
                    config = RunnableConfig(tags=["translation"])
                    res = await chain.ainvoke(
                        {"target_language": state['target_language'], "text": text},
                        config=config
                    )
                    return (text, res.content)
                except Exception as e:
                    if attempt < MAX_RETRIES:
                        wait_time = (2 ** attempt) * 0.5  # 指数退避
                        logger.warning(f"⚠️  重试 {attempt + 1}/{MAX_RETRIES}: {text[:20]}... ({e})")
                        await asyncio.sleep(wait_time)
                    else:
                        logger.error(f"❌ 最终失败: {text[:20]}... ({e})")
                        return (text, None)
    
    # 分批处理
    batches = [batch_texts[i:i + BATCH_SIZE] for i in range(0, len(batch_texts), BATCH_SIZE)]
    total_batches = len(batches)
    
    logger.info(f"📦 总计 {len(batch_texts)} 个文本，分成 {total_batches} 个批次处理")
    
    start_time = time.time()
    
    # 为每个批次创建任务
    batch_tasks = []
    for batch_idx, batch in enumerate(batches):
        # 创建批次内的所有翻译任务
        batch_length = len(batch)
        tasks = [translate_single(text) for text, _ in batch]
        
        # 创建批次处理任务（收集该批次的结果）
        async def process_batch(batch_idx: int, tasks: List, batch_length: int) -> None:
            logger.info(f"🚀 开始处理批次 {batch_idx + 1}/{total_batches}")
            results = await asyncio.gather(*tasks)
            
            # 处理批次结果
            batch_success = 0
            for original_text, translated_text in results:
                if translated_text:
                    translation_map[original_text] = translated_text
                    batch_success += 1
                else:
                    translation_map[original_text] = original_text
            
            logger.info(f"✅ 批次 {batch_idx + 1} 完成 ({batch_success}/{batch_length} 成功)")
        
        batch_tasks.append(process_batch(batch_idx, tasks, batch_length))
    
    # 并发执行所有批次
    await asyncio.gather(*batch_tasks)
    
    elapsed_time = time.time() - start_time
    
    logger.info(f"🎉 所有翻译完成！")
    logger.info(f"⏱️  总耗时: {elapsed_time:.2f} 秒")
    logger.info(f"🚀 平均每个文本: {elapsed_time/len(batch_texts):.2f} 秒")
    
    state["translation_map"] = translation_map
    state["status_msg"] = f"✅ 翻译完成，正在重构 PPT..."
    return state


def node_reconstruct_ppt(state: AgentState) -> AgentState:
    logger.info("🔨 开始智能重构 PPT ...")
    
    prs = Presentation(state['input_ppt_path'])
    translation_map = state["translation_map"]
    
    # 样式保持的关键参数
    MIN_FONT_SIZE_PT = 12
    MAX_FONT_REDUCTION = 0.5
    WIDTH_EXPANSION_LIMIT = 1.15
    MIN_RIGHT_MARGIN = Inches(0.3)
    MIN_LEFT_MARGIN = Inches(0.3)
    
    slide_width = prs.slide_width
    slide_height = prs.slide_height

    # 扩展宽度函数
    def expand_box_width_aware(
        shape, 
        alignment: PP_ALIGN,
        current_slide_boxes: List[Dict]
    ) -> bool:
        """
        根据对齐方式智能扩展文本框宽度 (修复方向性和边距问题)
        """
        old_width = shape.width
        old_left = shape.left
        old_right = old_left + old_width
        
        # 计算可用空间
        if alignment == PP_ALIGN.LEFT:
            # 左对齐：只能向右扩展
            max_possible_width = slide_width - MIN_RIGHT_MARGIN - old_left
            if max_possible_width <= old_width:
                return False
            
        elif alignment == PP_ALIGN.RIGHT:
            # 右对齐：只能向左扩展
            max_possible_width = old_right - MIN_LEFT_MARGIN
            if max_possible_width <= old_width:
                return False
            
        elif alignment == PP_ALIGN.CENTER:
            # 居中对齐：向两边扩展
            center = old_left + old_width / 2
            left_space = center - MIN_LEFT_MARGIN
            right_space = (slide_width - MIN_RIGHT_MARGIN) - center
            half_expansion = min(left_space, right_space)
            max_possible_width = half_expansion * 2
            
            if max_possible_width <= old_width:
                return False
            
            # 智能侧边扩展逻辑：如果一侧受阻，优先使用另一侧
            # 这里计算纯几何空间，后续碰撞检测会处理具体阻挡
            # 尝试非对称扩展的简单策略
            if left_space > right_space * 1.5:
                # 左侧空间大得多，尝试向左多扩一点（保持视觉中心感）
                # 这里暂不改变 center，仅在碰撞检测时微调
                pass
            elif right_space > left_space * 1.5:
                pass
                
        else:
            # 默认：左对齐处理
            max_possible_width = slide_width - MIN_RIGHT_MARGIN - old_left
            if max_possible_width <= old_width:
                return False
        
        # 计算目标宽度
        target_width = min(old_width * WIDTH_EXPANSION_LIMIT, max_possible_width)
        
        # 确保宽度增加（避免浮点误差）
        if target_width <= old_width:
            return False
        
        # 转换为整数
        target_width = int(target_width)
        
        # 预计算新的位置和尺寸
        new_left = old_left
        new_width = target_width
        
        if alignment == PP_ALIGN.CENTER:
            new_left = center - target_width / 2
            # 边界修正
            if new_left < MIN_LEFT_MARGIN:
                new_left = MIN_LEFT_MARGIN
                new_width = min(target_width, (center + old_width / 2) - MIN_LEFT_MARGIN)
            if new_left + new_width > slide_width - MIN_RIGHT_MARGIN:
                new_width = slide_width - MIN_RIGHT_MARGIN - new_left
                new_left = center - new_width / 2
        elif alignment == PP_ALIGN.RIGHT:
            new_left = old_right - new_width
            if new_left < MIN_LEFT_MARGIN:
                new_width = old_right - MIN_LEFT_MARGIN
                new_left = MIN_LEFT_MARGIN
        
        # 边界检查，防止负数宽度
        if new_width <= old_width:
            return False

        # 碰撞检测 (方向性过滤 + 零边距)
        test_box = {
            'left': new_left,
            'top': shape.top,
            'width': new_width,
            'height': shape.height,
            'shape_id': id(shape)
        }
        
        blocked_by = None
        
        for other_box in current_slide_boxes:
            if other_box['shape_id'] == id(shape):
                continue
            
            # --- 方向性过滤 ---
            # 1. 左对齐扩展向右：忽略完全在当前文本框左侧的物体
            if alignment == PP_ALIGN.LEFT:
                other_right = other_box['left'] + other_box['width']
                # 如果邻居在旧右边界的左侧，忽略它（我们在往右走）
                if other_right <= old_right + Inches(0.01):
                    continue
            
            # 2. 右对齐扩展向左：忽略完全在右侧的物体
            elif alignment == PP_ALIGN.RIGHT:
                # 如果邻居在旧左边界的右侧，忽略它
                if other_box['left'] >= old_left - Inches(0.01):
                    continue
            
            # 3. 居中对齐：两端都要检测，暂不做特殊过滤
            
            # 执行碰撞检测 (margin设为0，允许紧贴)
            if is_overlap(test_box, other_box, margin=Inches(0.0)):
                blocked_by = other_box['shape_id']
                # 如果被挡住，尝试回退
                break
        
        if blocked_by:
            # 如果居中对齐被挡，尝试偏移中心点（简单的挽救措施）
            if alignment == PP_ALIGN.CENTER:
                # 尝试只向没有阻挡的一侧扩展
                # 这里为了简化，如果居中被挡，直接返回失败
                # 因为偏移中心点会改变设计意图
                pass 
            return False
        
        # 应用修改
        shape.left = int(new_left)
        shape.width = int(new_width)
        return True

    # ========== 主处理流程 ==========
    
    # 收集文本框位置信息
    all_text_boxes = []
    for slide_idx, slide in enumerate(prs.slides):
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text.strip():
                all_text_boxes.append({
                    'slide_idx': slide_idx,
                    'shape_id': id(shape),
                    'left': shape.left,
                    'top': shape.top,
                    'width': shape.width,
                    'height': shape.height
                })
    
    replaced_count = 0
    adjustment_count = 0
    stats = {
        'font_reduced': 0,
        'width_expanded': 0,
        'wrap_enabled': 0,
        'no_adjustment': 0
    }
    
    for slide_idx, slide in enumerate(prs.slides):
        current_slide_boxes = [
            box for box in all_text_boxes 
            if box['slide_idx'] == slide_idx
        ]
        
        # 第一阶段：收集本页需要翻译的文本框信息
        group_candidates = defaultdict(list)
        
        for shape in slide.shapes:
            if not shape.has_text_frame:
                continue
            
            original_text = shape.text.strip()
            if original_text not in translation_map:
                continue
            
            translated_text = translation_map[original_text]
            
            # 获取特征属性用于分组
            font_size = get_font_size(shape)
            alignment = get_paragraph_alignment(shape)
            font_name = "Arial"
            for p in shape.text_frame.paragraphs:
                for r in p.runs:
                    if r.font.name:
                        font_name = r.font.name
                        break
            
            # 计算 ratio
            length_ratio = get_visual_width_ratio(original_text, translated_text)
            
            group_candidates[(font_size.pt, alignment, font_name)].append({
                'shape': shape,
                'original_text': original_text,
                'translated_text': translated_text,
                'font_size_pt': font_size.pt,
                'length_ratio': length_ratio,
                'has_numbers': has_arabic_numbers(translated_text)
            })
        
        # 第二阶段：按组处理
        for group_key, group_members in group_candidates.items():
            base_size, base_align, base_font = group_key
            member_count = len(group_members)
            if member_count == 0: continue
            
            # 计算全组的统一调整策略
            ratios = [m['length_ratio'] for m in group_members]
            max_ratio = max(ratios)
            median_ratio = median(ratios)
            
            # 抗干扰算法
            if max_ratio > 2.5 or max_ratio > median_ratio * 1.5:
                effective_ratio = min(median_ratio * 1.5, 2.5)
                logger.info(f"  📦 组处理 (字号={base_size}pt): 检测到异常值 (Max={max_ratio:.2f}, Median={median_ratio:.2f}), 采用比例上限 {effective_ratio:.2f}")
            else:
                effective_ratio = max_ratio
                logger.info(f"  📦 组处理 (字号={base_size}pt): 正常调整 (Max={max_ratio:.2f})")
            
            reduction_ratio = calculate_dynamic_reduction_ratio(effective_ratio)
            
            # 应用样式替换和同步字号
            for member in group_members:
                shape = member['shape']
                original_text = member['original_text']
                translated_text = member['translated_text']

                original_top = shape.top
                original_left = shape.left
                original_width = shape.width
                original_height = shape.height

                # A. 保存样式
                original_styles = []
    
                for para_idx, paragraph in enumerate(shape.text_frame.paragraphs):
                    para_alignment = paragraph.alignment
                    para_space_before = paragraph.space_before
                    para_space_after = paragraph.space_after
                    para_level = paragraph.level
                    
                    p_elem = paragraph._p
                    bullet_info = extract_bullet_info_from_xml(p_elem)
                    
                    for run_idx, run in enumerate(paragraph.runs):
                        if not run.text:
                            continue
                        
                        # 确保所有bullet相关字段都存在（避免KeyError）
                        style = {
                            'paragraph_idx': para_idx,
                            'run_idx': run_idx,
                            'text': run.text,
                            # 段落级
                            'alignment': para_alignment,
                            'space_before': para_space_before,
                            'space_after': para_space_after,
                            'level': para_level,
                            
                            # ===== 项目符号/编号信息（确保所有字段都有默认值）=====
                            'has_bullet': bullet_info.get('has_bullet', False),
                            'bullet_type': bullet_info.get('bullet_type', 'inherited'),
                            
                            # 项目符号专用（即使不是char类型也设置None）
                            'bullet_char': bullet_info.get('bullet_char', None),
                            
                            # 编号专用（即使不是autoNum类型也设置None）
                            'auto_num_type': bullet_info.get('auto_num_type', None),
                            'auto_num_start': bullet_info.get('auto_num_start', 1),
                            
                            # 共用样式（所有类型都可能有）
                            'bullet_font_name': bullet_info.get('bullet_font_name', None),
                            'bullet_font_size': bullet_info.get('bullet_font_size', None),
                            'bullet_color': bullet_info.get('bullet_color', None),
                            'bullet_color_type': bullet_info.get('bullet_color_type', None),
                            'bullet_level': bullet_info.get('level', 0),
                            'bullet_marL': bullet_info.get('marL', None),
                            'bullet_indent': bullet_info.get('indent', None),
                            
                            # 字符级
                            'font_name': run.font.name,
                            'font_size': run.font.size,
                            'font_bold': run.font.bold,
                            'font_italic': run.font.italic,
                            'font_underline': run.font.underline,
                            'color': None,
                            'color_type': None,
                        }
                        
                        # 提取字符颜色
                        if run.font.color:
                            if hasattr(run.font.color, 'rgb') and run.font.color.rgb:
                                style['color'] = run.font.color.rgb
                                style['color_type'] = 'RGB'
                            elif hasattr(run.font.color, 'theme_color') and run.font.color.theme_color:
                                style['color'] = run.font.color.theme_color
                                style['color_type'] = 'theme'
                        
                        original_styles.append(style)
                
                # B. 替换文字
                shape.text = translated_text
                shape.top = original_top
                shape.left = original_left
                shape.width = original_width
                shape.height = original_height
                
                # C. 恢复样式
                apply_styles(shape, original_styles)
                
                # D. 设置自动调整选项 (不改变形状大小，允许文本溢出)
                text_frame = shape.text_frame
                try:
                    text_frame.auto_size = MSO_AUTO_SIZE.NONE
                    # text_frame.word_wrap = False
                except Exception as e:
                    logger.warning(f"设置 auto_size 失败: {e}")

                for para_idx, paragraph in enumerate(text_frame.paragraphs):
                    # 重新从 original_styles 中获取对齐方式并应用
                    for style in original_styles:
                        if style.get('paragraph_idx') == para_idx and style.get('alignment'):
                            paragraph.alignment = style['alignment']
                            break
                
                # E. 应用组统一的字号
                new_font_size_pt = base_size * reduction_ratio
                new_font_size_pt = max(new_font_size_pt, MIN_FONT_SIZE_PT)
                new_font_size = Pt(new_font_size_pt)
                for paragraph in text_frame.paragraphs:
                    for run in paragraph.runs:
                        if run.font.size:
                            run.font.size = new_font_size
                shape.top = original_top
                shape.left = original_left
                # width 和 height 可能已被 expand_box_width_aware 修改，所以只在需要时恢复
                if not (new_font_size_pt <= MIN_FONT_SIZE_PT and member['length_ratio'] > 1.2) and member['length_ratio'] <= 2.0:
                    # 只有在不扩展宽度的情况下，才恢复原始宽度和高度
                    shape.width = original_width
                    shape.height = original_height

                replaced_count += 1
                
                # 第三阶段：个别优化
                real_ratio = member['length_ratio']
                
                is_overcrowded = (new_font_size_pt <= MIN_FONT_SIZE_PT and real_ratio > 1.2) or (real_ratio > 2.0)
                
                if is_overcrowded:
                    # 调用修复后的函数
                    success = expand_box_width_aware(
                        shape,
                        base_align,
                        current_slide_boxes
                    )
                    
                    if success:
                        stats['width_expanded'] += 1
                        adjustment_count += 1
                        logger.info(f"    ↔️  扩展宽度成功: {translated_text[:15]}...")
                        # 宽度失败，尝试换行
                    try:
                        text_frame.word_wrap = True
                        stats['wrap_enabled'] += 1
                        logger.info(f"启用换行: {translated_text[:15]}...")
                    except: pass
                elif effective_ratio <= 1.05:
                    stats['no_adjustment'] += 1
                else:
                    stats['font_reduced'] += 1
                    adjustment_count += 1
                    logger.info(f"    📏 同步字号: {base_size}pt -> {new_font_size_pt}pt")

    # 保存文件
    output_ppt_path = state.get('output_ppt_path')
    if not output_ppt_path:
        input_ppt_path = state.get('input_ppt_path')
        target_lang = state.get('target_language')
        path = Path(input_ppt_path)
        new_filename = f"{path.stem}_{target_lang}{path.suffix}"
        output_ppt_path = str(path.parent / new_filename)
        
    prs.save(output_ppt_path)
    
    # 输出统计信息
    logger.info(f"✅ 重构完成！")
    logger.info(f"   - 共替换 {replaced_count} 处文本")
    logger.info(f"   - 总计调整 {adjustment_count} 处")
    logger.info(f"   ├─ 缩小字号: {stats['font_reduced']}")
    logger.info(f"   ├─ 扩展宽度: {stats['width_expanded']}")
    logger.info(f"   ├─ 启用换行: {stats['wrap_enabled']}")
    logger.info(f"   └─ 无需调整: {stats['no_adjustment']}")
    
    state["status_msg"] = f"✅ PPT 生成成功！共翻译 {replaced_count} 处，调整 {adjustment_count} 处"
    return state

# ==========================================
# 4. 包装异步节点以适配 LangGraph
# ==========================================
def make_translate_node(llm):
    """
    工厂函数：接收llm，返回绑定了llm的同步包装函数（函数对象）
    作用：让wrapper_translate_text能拿到llm，且返回的是可调用的函数对象
    """
    def wrapper_translate_text(state: AgentState) -> AgentState:
        """包装异步节点为同步函数（现在绑定了llm）"""
        loop = asyncio.new_event_loop()
        asyncio.set_event_loop(loop)
        try:
            # 调用异步节点时传入绑定的llm
            return loop.run_until_complete(async_node_translate_text(llm, state))
        finally:
            loop.close()
    return wrapper_translate_text

# ==========================================
# 4. 构建 LangGraph 工作流
# ==========================================
def create_graph(llm):
    workflow = StateGraph(AgentState)
    
    workflow.add_node("parse", node_parse_ppt)
    workflow.add_node("translate", make_translate_node(llm))  # 使用包装后的同步节点
    workflow.add_node("reconstruct", node_reconstruct_ppt)
    
    workflow.set_entry_point("parse")
    workflow.add_edge("parse", "translate")
    workflow.add_edge("translate", "reconstruct")
    workflow.add_edge("reconstruct", END)
    
    return workflow.compile()