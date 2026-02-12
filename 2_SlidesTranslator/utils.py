import os
import re
from typing import List, Dict, Optional
from pptx.util import Pt, Inches, Emu
from pptx.oxml.ns import qn
from lxml import etree
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN, MSO_ANCHOR
from collections import Counter


# 加载提示词
def load_prompt(file_path: str) -> str:
    """
    读取 prompt 文本文件
    """
    if not os.path.exists(file_path):
        raise FileNotFoundError(f"Prompt 文件不存在: {file_path}")
        
    with open(file_path, "r", encoding="utf-8") as f:
        return f.read()

# 视觉宽度修正计算
def get_visual_width_ratio(original_text: str, translated_text: str) -> float:
        """
        计算文本的视觉宽度比例
        
        视觉宽度权重系数：
        - CJK (中日韩): 1.85 (1个汉字 ≈ 1.85个英文字母宽度)
        - Latin (英文等): 1.0 (基准)
        - Mixed: 根据内容混合计算
        """
        WIDTH_WEIGHTS = {
            'CJK': 1.85,   # 中日韩
            'Latin': 1.0   # 英文、数字等
        }
        
        def detect_script_type(text: str) -> str:
            """
            检测文本的主语系
            如果 CJK 字符占比超过 20%，则视为 CJK 文本
            """
            # CJK 字符正则：包括中文、日文、韩文
            cjk_chars = re.findall(r'[\u4e00-\u9fff\u3040-\u309f\u30a0-\u30ff\uac00-\ud7af]', text)
            total_chars = len([c for c in text if c.strip()])
            
            if not total_chars:
                return 'Latin'
                
            cjk_ratio = len(cjk_chars) / total_chars
            return 'CJK' if cjk_ratio > 0.2 else 'Latin'

        orig_type = detect_script_type(original_text)
        trans_type = detect_script_type(translated_text)
        
        orig_visual_len = len(original_text) * WIDTH_WEIGHTS[orig_type]
        trans_visual_len = len(translated_text) * WIDTH_WEIGHTS[trans_type]
        
        if orig_visual_len == 0:
            return 1.0
        
        ratio = trans_visual_len / orig_visual_len
        return ratio

# 文本框是否重叠检测
def is_overlap(box1: Dict, box2: Dict, margin: float = Inches(0.05)) -> bool:
        left1, top1, right1, bottom1 = (
            int(box1['left'] - margin), int(box1['top'] - margin),
            int(box1['left'] + box1['width'] + margin), int(box1['top'] + box1['height'] + margin)
        )
        left2, top2, right2, bottom2 = (
            int(box2['left'] - margin), int(box2['top'] - margin),
            int(box2['left'] + box2['width'] + margin), int(box2['top'] + box2['height'] + margin)
        )
        not_overlap = (right1 < left2 or right2 < left1 or bottom1 < top2 or bottom2 < top1)
        return not not_overlap

# 检测是否存在阿拉伯数字
def has_arabic_numbers(text: str) -> bool:
        return bool(re.search(r'\d+', text))

# 获取文本框中文字的平均字号
def get_font_size(shape) -> Optional[Pt]:
    sizes = []
    for paragraph in shape.text_frame.paragraphs:
        for run in paragraph.runs:
            if run.font.size:
                sizes.append(run.font.size.pt)
    return Pt(sum(sizes) / len(sizes)) if sizes else Pt(18)

# 获取文本框的主要对齐方式 
def get_paragraph_alignment(shape) -> Optional[PP_ALIGN]:
    alignments = []
    for paragraph in shape.text_frame.paragraphs:
        if paragraph.alignment:
            alignments.append(paragraph.alignment)
    
    if not alignments:
        return PP_ALIGN.LEFT
    
    return Counter(alignments).most_common(1)[0][0]

# 将原有样式应用到修改后的文本框
def apply_styles(shape, original_styles: List[Dict]) -> None:
    text_frame = shape.text_frame
    style_idx = 0

    for para_idx, paragraph in enumerate(text_frame.paragraphs):
        if style_idx >= len(original_styles):
            break
            
        style = original_styles[style_idx]
        p_elem = paragraph._p

        # --- 段落级格式 ---
        if style.get('alignment'):
            paragraph.alignment = style['alignment']
        if style.get('space_before') is not None:
            paragraph.space_before = style['space_before']
        if style.get('space_after') is not None:
            paragraph.space_after = style['space_after']
        if style.get('level') is not None:
            paragraph.level = style['level']

        # --- 项目符号/编号样式 ---
        apply_bullet_style(p_elem, style)

        # --- 字符级格式 ---
        for run in paragraph.runs:
            if style.get('font_name'):
                run.font.name = style['font_name']
            if style.get('font_size'):
                run.font.size = style['font_size']
            if style.get('font_bold') is not None:
                run.font.bold = style['font_bold']
            if style.get('font_italic') is not None:
                run.font.italic = style['font_italic']
            if style.get('font_underline') is not None:
                run.font.underline = style['font_underline']
            
            # 颜色处理
            if style.get('color_type') == 'RGB' and style.get('color'):
                try:
                    run.font.color.rgb = style['color']
                except Exception:
                    pass
            elif style.get('color_type') == 'theme' and style.get('color'):
                try:
                    run.font.color.theme_color = style['color']
                except Exception:
                    pass

        style_idx += 1

# 根据 length_ratio 动态计算字号缩小比例
def calculate_dynamic_reduction_ratio(length_ratio: float) -> float:
    """
    根据翻译后/原文本长度比，**平滑动态**计算字号缩小比例
    """
    if not isinstance(length_ratio, (int, float)) or length_ratio <= 0:
        return 1.0

    RATIO_ANCHORS = {
        1.2: 0.95,
        1.5: 0.85,
        2.0: 0.80,  
        2.5: 0.70,  
        3.0: 0.60,  
        4.0: 0.50   
    }
    
    ratio_list = sorted(RATIO_ANCHORS.keys())
    reduction_list = [RATIO_ANCHORS[k] for k in ratio_list]

    if length_ratio <= 1.0:
        return 1.0

    if length_ratio in RATIO_ANCHORS:
        return RATIO_ANCHORS[length_ratio]

    # 线性插值
    for i in range(len(ratio_list) - 1):
        r1, r2 = ratio_list[i], ratio_list[i+1]
        red1, red2 = reduction_list[i], reduction_list[i+1]
        if r1 < length_ratio < r2:
            interpolation = red1 + (red2 - red1) * (length_ratio - r1) / (r2 - r1)
            return round(interpolation, 3)

    # 超过最大锚点后缓慢缩小
    max_ratio = ratio_list[-1]  # 4.0
    max_red = reduction_list[-1]  # 0.50
    min_red = 0.50  # 最小 50%
    slope = -0.01
    extra_red = max_red + slope * (length_ratio - max_ratio)
    return round(max(extra_red, min_red), 3)

# 从段落XML元素中提取完整的项目符号/编号信息
def extract_bullet_info_from_xml(p_elem) -> Dict:
    bullet_info = {
        'has_bullet': False,
        'bullet_type': 'inherited',  # 默认假设继承
        # ... 其他字段
    }
    
    # 查找段落属性
    pPr = None
    for child in p_elem:
        if child.tag == '{http://schemas.openxmlformats.org/drawingml/2006/main}pPr':
            pPr = child
            break
    
    # 没有 pPr = 肯定是继承
    if pPr is None:
        return bullet_info  # bullet_type = 'inherited'
    
    # 获取级别和缩进
    bullet_info['level'] = int(pPr.get('lvl', '0'))
    bullet_info['marL'] = pPr.get('marL')
    bullet_info['indent'] = pPr.get('indent')
    
    # 检查是否有显式的项目符号标记
    found_explicit_bullet = False
    
    for child in pPr:
        tag = child.tag.split('}')[-1]
        
        if tag == 'buNone':
            bullet_info['has_bullet'] = False
            bullet_info['bullet_type'] = 'none'
            found_explicit_bullet = True
            return bullet_info  # 明确无项目符号
        
        elif tag == 'buChar':
            bullet_info['has_bullet'] = True
            bullet_info['bullet_type'] = 'char'
            bullet_info['bullet_char'] = child.get('char', '•')
            found_explicit_bullet = True
        
        elif tag == 'buAutoNum':
            bullet_info['has_bullet'] = True
            bullet_info['bullet_type'] = 'autoNum'
            bullet_info['auto_num_type'] = child.get('type', 'arabicPlain')
            bullet_info['auto_num_start'] = int(child.get('startAt', '1'))
            found_explicit_bullet = True
        
        elif tag == 'buBlip':
            bullet_info['has_bullet'] = True
            bullet_info['bullet_type'] = 'blip'
            found_explicit_bullet = True
        
        # 提取样式属性（buFont, buSzPct, buClr）
        elif tag == 'buFont':
            bullet_info['bullet_font_name'] = child.get('typeface')
        elif tag == 'buSzPct':
            val = child.get('val', '100000')
            bullet_info['bullet_font_size'] = int(val) / 1000
        elif tag == 'buClr':
            # 提取颜色...
            pass
    
    # 关键判断：只有找到显式标记时才不是继承
    if found_explicit_bullet:
        return bullet_info
    else:
        # 有 pPr 但没有 bu* 元素 = 继承
        bullet_info['bullet_type'] = 'inherited'
        bullet_info['has_bullet'] = 'unknown'  # 未知，取决于母版
        return bullet_info

# 应用项目符号/编号样式
def apply_bullet_style(p_elem, style: Dict) -> None:
    # 获取或创建 <a:pPr>
    pPr = None
    for child in p_elem:
        if child.tag == qn('a:pPr'):
            pPr = child
            break
    
    if pPr is None:
        pPr = etree.SubElement(p_elem, qn('a:pPr'))
    
    # 设置级别和缩进
    if style.get('bullet_level') is not None:
        pPr.set('lvl', str(style['bullet_level']))
    if style.get('bullet_marL') is not None:
        pPr.set('marL', str(style['bullet_marL']))
    if style.get('bullet_indent') is not None:
        pPr.set('indent', str(style['bullet_indent']))
    
    bullet_type = style.get('bullet_type')
    
    # 继承模式：不清理XML
    if bullet_type == 'inherited':
        return
    
    # 清除现有的项目符号/编号子元素
    bullet_tags = ['buNone', 'buChar', 'buAutoNum', 'buBlip', 'buFont', 'buSzPct', 'buClr']
    for child in list(pPr):
        tag = child.tag.split('}')[-1]
        if tag in bullet_tags:
            pPr.remove(child)
    
    # === 1. 无项目符号 ===
    if bullet_type == 'none':
        etree.SubElement(pPr, qn('a:buNone'))
        return
    
    # === 2. 自定义字符项目符号 ===
    elif bullet_type == 'char':
        # 修复：使用.get()提供默认值，避免KeyError
        bullet_char = style.get('bullet_char', '•')  # 默认黑点
        if not bullet_char:  # 如果是空字符串也使用默认
            bullet_char = '•'
        
        buChar = etree.SubElement(pPr, qn('a:buChar'))
        buChar.set('char', bullet_char)
        _add_bullet_style_elements(pPr, style)
    
    # === 3. 自动编号 ===
    elif bullet_type == 'autoNum':
        # 修复：使用.get()提供默认值
        auto_num_type = style.get('auto_num_type', 'arabicPlain')
        auto_num_start = style.get('auto_num_start', 1)
        
        buAutoNum = etree.SubElement(pPr, qn('a:buAutoNum'))
        buAutoNum.set('type', auto_num_type)
        buAutoNum.set('startAt', str(auto_num_start))
        _add_bullet_style_elements(pPr, style)
    
    # === 4. 图片项目符号 ===
    elif bullet_type == 'blip':
        buBlip = etree.SubElement(pPr, qn('a:buBlip'))
    
    # === 5. 未知类型，视为继承 ===
    else:
        # 如果bullet_type是None或其他未知值，不要清理XML
        pass

# 辅助函数：添加项目符号/编号的共用样式元素
def _add_bullet_style_elements(pPr, style: Dict) -> None:
    # 字体
    bullet_font = style.get('bullet_font_name')
    if bullet_font and bullet_font != 'default':
        buFont = etree.SubElement(pPr, qn('a:buFont'))
        buFont.set('typeface', bullet_font)
        buFont.set('pitchFamily', '34')
        buFont.set('charset', '0')
    
    # 大小
    bullet_size = style.get('bullet_font_size')
    if bullet_size is not None:
        buSzPct = etree.SubElement(pPr, qn('a:buSzPct'))
        buSzPct.set('val', str(int(bullet_size * 1000)))
    
    # 颜色
    bullet_color = style.get('bullet_color')
    bullet_color_type = style.get('bullet_color_type')
    if bullet_color and bullet_color_type:
        buClr = etree.SubElement(pPr, qn('a:buClr'))
        if bullet_color_type == 'RGB':
            srgbClr = etree.SubElement(buClr, qn('a:srgbClr'))
            # 修复：处理不同类型的color值
            if isinstance(bullet_color, str):
                srgbClr.set('val', bullet_color.replace('#', ''))
            else:
                srgbClr.set('val', str(bullet_color))
        elif bullet_color_type == 'theme':
            schemeClr = etree.SubElement(buClr, qn('a:schemeClr'))
            schemeClr.set('val', str(bullet_color))

# 常见的自动编号类型映射
AUTO_NUM_TYPES = {
    # 阿拉伯数字
    'arabicPlain': '1, 2, 3...',
    'arabicPeriod': '1., 2., 3....',
    'arabicParenR': '1), 2), 3)...',
    'arabicParenBoth': '(1), (2), (3)...',
    
    # 小写字母
    'alphaLcPlain': 'a, b, c...',
    'alphaLcPeriod': 'a., b., c....',
    'alphaLcParenR': 'a), b), c)...',
    'alphaLcParenBoth': '(a), (b), (c)...',
    
    # 大写字母
    'alphaUcPlain': 'A, B, C...',
    'alphaUcPeriod': 'A., B., C....',
    'alphaUcParenR': 'A), B), C)...',
    'alphaUcParenBoth': '(A), (B), (C)...',
    
    # 小写罗马数字
    'romanLcPlain': 'i, ii, iii...',
    'romanLcPeriod': 'i., ii., iii....',
    'romanLcParenR': 'i), ii), iii)...',
    'romanLcParenBoth': '(i), (ii), (iii)...',
    
    # 大写罗马数字
    'romanUcPlain': 'I, II, III...',
    'romanUcPeriod': 'I., II., III....',
    'romanUcParenR': 'I), II), III)...',
    'romanUcParenBoth': '(I), (II), (III)...',
    
    # 特殊符号
    'circleNumWdBlackPlain': '❶, ❷, ❸...',
    'circleNumWdWhitePlain': '①, ②, ③...',
}